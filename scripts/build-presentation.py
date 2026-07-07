#!/usr/bin/env python3
"""Build a sexy PPTX presentation from scraped TikTok AI trends data."""

import json
import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
DATA_FILE = os.path.join(PROJECT_DIR, "outputs", "tiktok-ai-raw.json")
OUTPUT_FILE = os.path.join(PROJECT_DIR, "outputs", "TikTok_AI_Trends_March2026.pptx")
THUMB_DIR = os.path.join(PROJECT_DIR, "outputs", "thumbs")

# Brand colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x0F)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_PINK = RGBColor(0xFE, 0x2C, 0x55)  # TikTok red/pink
ACCENT_CYAN = RGBColor(0x25, 0xF4, 0xEE)  # TikTok cyan
ACCENT_PURPLE = RGBColor(0xA855, 0xF7, 0x00)[0:3] if False else RGBColor(0xA8, 0x55, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x10, 0xB9, 0x81)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def download_thumbnail(url, filename):
    """Download a thumbnail image, return path or None."""
    os.makedirs(THUMB_DIR, exist_ok=True)
    path = os.path.join(THUMB_DIR, filename)
    if os.path.exists(path):
        return path
    try:
        headers = {"User-Agent": "Mozilla/5.0", "Referer": "https://www.tiktok.com/"}
        r = requests.get(url, timeout=15, headers=headers)
        if r.status_code == 200 and len(r.content) > 1000:
            with open(path, "wb") as f:
                f.write(r.content)
            return path
    except Exception as e:
        print(f"  Failed to download {filename}: {e}")
    return None


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines):
    """Add text box with multiple styled lines. lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Arial"
        p.space_after = Pt(4)
    return txBox


def fmt_num(n):
    """Format number with K/M suffix."""
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n/1_000:.0f}K"
    return str(n)


def build_title_slide(prs, videos):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # TikTok-style gradient accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_PINK
    bar.line.fill.background()

    # Cyan accent bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.08), SLIDE_WIDTH, Inches(0.04))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_CYAN
    bar2.line.fill.background()

    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "TRENDING TIKTOK TOPICS", 48, ACCENT_CYAN, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "IN THE AI NICHE", 56, WHITE, True, PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
                 "March 2026  |  Live data scraped via Apify", 22, GRAY, False, PP_ALIGN.CENTER)

    # Stats row
    stats = [
        ("47", "Videos Scraped"),
        ("15", "Transcripts"),
        ("62.8M", "Top Views"),
        ("4.5%", "Avg Engagement"),
    ]
    x_start = Inches(1.5)
    for i, (num, label) in enumerate(stats):
        x = x_start + Inches(2.7) * i
        card = add_shape(slide, x, Inches(4.8), Inches(2.2), Inches(1.6), BG_CARD, 0.05)
        add_text_box(slide, x + Inches(0.1), Inches(4.9), Inches(2), Inches(0.9),
                     num, 36, ACCENT_PINK, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(5.6), Inches(2), Inches(0.5),
                     label, 14, GRAY, False, PP_ALIGN.CENTER)

    # Bottom bar
    bar3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.38), SLIDE_WIDTH, Inches(0.12))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = ACCENT_PINK
    bar3.line.fill.background()


def build_categories_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "6 TRENDING CONTENT CATEGORIES", 36, ACCENT_CYAN, True, PP_ALIGN.LEFT)

    categories = [
        ("01", "AI Entertainment", "62.8M top", "Funny AI-generated videos, dancing, face swaps, talking objects", ACCENT_PINK),
        ("02", "AI Tools Listicles", "1.7M top", '"Use this for X" rapid-fire roundups of 6-8 tools', ACCENT_CYAN),
        ("03", "AI Industry Drama", "437K top", "Anthropic vs OpenAI, Pentagon controversy, CEO leaks", ACCENT_PURPLE),
        ("04", "AI Tutorials", "286K top", "Step-by-step creation guides with screen recordings", GREEN),
        ("05", "AI Tips & Secrets", "185K top", '"Things you didn\'t know about Claude/ChatGPT"', RGBColor(0xFF, 0xD7, 0x00)),
        ("06", "AI Business/Money", "401K top", "AI agents with crypto wallets, $10K of work for $32/mo", RGBColor(0xFF, 0x69, 0x34)),
    ]

    for i, (num, title, views, desc, color) in enumerate(categories):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.3) + row * Inches(2.0)

        card = add_shape(slide, x, y, Inches(5.9), Inches(1.8), BG_CARD, 0.03)

        # Number accent
        num_shape = add_shape(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.7), Inches(0.7), color, 0.1)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.35), Inches(0.7), Inches(0.5),
                     num, 22, WHITE, True, PP_ALIGN.CENTER)

        # Title + views
        add_text_box(slide, x + Inches(1.1), y + Inches(0.25), Inches(3.5), Inches(0.5),
                     title, 22, WHITE, True)
        add_text_box(slide, x + Inches(4.6), y + Inches(0.3), Inches(1.1), Inches(0.4),
                     views, 14, color, True, PP_ALIGN.RIGHT)

        # Description
        add_text_box(slide, x + Inches(1.1), y + Inches(0.8), Inches(4.5), Inches(0.8),
                     desc, 13, LIGHT_GRAY, False)


def build_top_videos_slide(prs, videos, start_idx, count, slide_title):
    """Build a slide showing video cards with thumbnails."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 slide_title, 32, ACCENT_CYAN, True, PP_ALIGN.LEFT)

    subset = videos[start_idx:start_idx + count]
    cols = min(count, 5)
    card_w = Inches(2.3)
    gap = Inches(0.2)
    total_w = cols * card_w + (cols - 1) * gap
    x_start = (SLIDE_WIDTH - total_w) / 2

    for i, v in enumerate(subset):
        row = i // 5
        col = i % 5
        x = x_start + col * (card_w + gap)
        y = Inches(1.2) + row * Inches(3.1)

        # Card background
        card = add_shape(slide, x, y, card_w, Inches(2.9), BG_CARD, 0.03)

        # Thumbnail
        cover_url = v.get("videoMeta", {}).get("coverUrl", "")
        thumb_path = None
        if cover_url:
            safe_name = f"thumb_{v.get('id', i)}.jpg"
            thumb_path = download_thumbnail(cover_url, safe_name)

        if thumb_path:
            try:
                pic = slide.shapes.add_picture(thumb_path, x + Inches(0.1), y + Inches(0.1),
                                               card_w - Inches(0.2), Inches(1.4))
            except Exception:
                pass

        # Author
        author = v.get("authorMeta", {}).get("name", "Unknown")
        add_text_box(slide, x + Inches(0.1), y + Inches(1.55), card_w - Inches(0.2), Inches(0.35),
                     f"@{author}", 12, ACCENT_CYAN, True)

        # Views + likes
        views = v.get("playCount", 0)
        likes = v.get("diggCount", 0)
        add_text_box(slide, x + Inches(0.1), y + Inches(1.85), card_w - Inches(0.2), Inches(0.3),
                     f"{fmt_num(views)} views  |  {fmt_num(likes)} likes", 11, ACCENT_PINK, True)

        # Caption snippet
        caption = v.get("text", "")[:80]
        if len(v.get("text", "")) > 80:
            caption += "..."
        add_text_box(slide, x + Inches(0.1), y + Inches(2.15), card_w - Inches(0.2), Inches(0.65),
                     caption, 9, GRAY, False)


def build_hooks_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "TOP PERFORMING HOOKS", 36, ACCENT_PINK, True, PP_ALIGN.LEFT)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(10), Inches(0.5),
                 "Extracted from video transcripts — what makes viewers stay", 16, GRAY, False)

    hooks = [
        ('"I finally stopped gatekeeping"', "62.8M", "Curiosity / FOMO", ACCENT_PINK),
        ('"Do you wanna [X]? Use this." (repeated)', "1.7M", "Rapid-fire value", ACCENT_CYAN),
        ('"Here\'s why everyone\'s canceling ChatGPT"', "437K", "Controversy", ACCENT_PURPLE),
        ('"13,000 AI Agents opened crypto wallets in a single day"', "401K", "Shocking stat", GREEN),
        ('"Things you didn\'t know you can do with Claude"', "185K", "Exclusivity", RGBColor(0xFF, 0xD7, 0x00)),
        ('"2.5M people have uninstalled ChatGPT, myself included"', "170K", "Social proof", RGBColor(0xFF, 0x69, 0x34)),
        ('"ChatGPT is fighting back today"', "143K", "Breaking news", ACCENT_CYAN),
        ('"Let me show you how to do it"', "117K", "Tutorial promise", GREEN),
    ]

    for i, (hook, views, hook_type, color) in enumerate(hooks):
        y = Inches(1.5) + i * Inches(0.72)

        # Accent bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.06), Inches(0.55))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        # Hook text
        add_text_box(slide, Inches(0.9), y, Inches(7.5), Inches(0.55),
                     hook, 16, WHITE, False)

        # Views badge
        badge = add_shape(slide, Inches(8.8), y + Inches(0.05), Inches(1.2), Inches(0.45), color, 0.15)
        add_text_box(slide, Inches(8.8), y + Inches(0.08), Inches(1.2), Inches(0.4),
                     views, 16, WHITE, True, PP_ALIGN.CENTER)

        # Type label
        add_text_box(slide, Inches(10.2), y + Inches(0.05), Inches(2.5), Inches(0.45),
                     hook_type, 13, GRAY, False)


def build_tools_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "MOST MENTIONED AI TOOLS", 36, ACCENT_CYAN, True, PP_ALIGN.LEFT)

    add_text_box(slide, Inches(0.6), Inches(0.9), Inches(10), Inches(0.5),
                 "Tools creators are recommending in viral AI TikToks", 16, GRAY, False)

    tools = [
        ("ChatGPT", "Save time, general AI assistant", "Most mentioned overall"),
        ("Claude", "Email writing, code, research", "Trending — Pentagon controversy boost"),
        ("Nano Banana Pro", "Ultra-realistic AI image generation", "Rising fast in tutorials"),
        ("Gemini AI", "Video analysis, image generation", "Google's multimodal play"),
        ("Higgsfield", "AI video creation + motion control", "New favorite for Reels"),
        ("RunwayML", "AI video editing & generation", "Established creative tool"),
        ("Bolt.new", "AI website builder", "Dev community favorite"),
        ("NotebookLM", "Learn new skills with AI", "Google's sleeper hit"),
        ("Replit", "Build apps with AI coding", "No-code revolution"),
        ("Perplexity", "AI-powered research", "Replacing Google Search"),
    ]

    for i, (name, use_case, note) in enumerate(tools):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.5) + row * Inches(1.1)

        card = add_shape(slide, x, y, Inches(5.9), Inches(0.95), BG_CARD, 0.03)

        # Tool name
        add_text_box(slide, x + Inches(0.3), y + Inches(0.1), Inches(2.5), Inches(0.4),
                     name, 18, ACCENT_PINK, True)

        # Use case
        add_text_box(slide, x + Inches(0.3), y + Inches(0.5), Inches(2.8), Inches(0.4),
                     use_case, 11, LIGHT_GRAY, False)

        # Note badge
        add_text_box(slide, x + Inches(3.2), y + Inches(0.3), Inches(2.5), Inches(0.4),
                     note, 10, ACCENT_CYAN, False, PP_ALIGN.RIGHT)


def build_takeaways_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Bottom accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.38), SLIDE_WIDTH, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_PINK
    bar.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "ACTIONABLE TAKEAWAYS", 36, ACCENT_PINK, True, PP_ALIGN.LEFT)

    takeaways = [
        ("Entertainment-first wins", "The highest-performing AI content is funny and surprising, not educational. Lead with entertainment, embed the AI angle.", ACCENT_PINK),
        ("Listicle format is reliable", '"Use this for X" with 6-8 tools consistently gets 100K+ views. Fast pace, 2-3 seconds per tool.', ACCENT_CYAN),
        ("AI drama is a goldmine", "Anthropic vs OpenAI controversy is driving massive engagement. Take a stance, create tribalism.", ACCENT_PURPLE),
        ("Use concrete numbers in hooks", '$10,000... 13,000 agents... 2.5M uninstalls. Specific numbers create irresistible curiosity gaps.', GREEN),
        ("Claude is trending hard", "Multiple viral videos about Claude features + Anthropic controversy boosting awareness. Good time to make Claude content.", RGBColor(0xFF, 0xD7, 0x00)),
        ("Tutorials need visual hooks", "Before/after or transformation reveals outperform dry screen recordings. Show the result first.", RGBColor(0xFF, 0x69, 0x34)),
    ]

    for i, (title, desc, color) in enumerate(takeaways):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.3) + row * Inches(2.0)

        card = add_shape(slide, x, y, Inches(5.9), Inches(1.8), BG_CARD, 0.03)

        # Color accent on left
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), Inches(1.8))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        # Number
        num_txt = f"0{i+1}"
        add_text_box(slide, x + Inches(0.25), y + Inches(0.15), Inches(0.6), Inches(0.5),
                     num_txt, 28, color, True)

        # Title
        add_text_box(slide, x + Inches(0.9), y + Inches(0.2), Inches(4.7), Inches(0.45),
                     title, 20, WHITE, True)

        # Description
        add_text_box(slide, x + Inches(0.9), y + Inches(0.7), Inches(4.7), Inches(0.9),
                     desc, 13, LIGHT_GRAY, False)


def build_hashtags_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 "HASHTAG PERFORMANCE", 36, ACCENT_CYAN, True, PP_ALIGN.LEFT)

    hashtags = [
        ("#ai", 10, "Broadest reach, highest competition", 100),
        ("#chatgpt", 10, "Product-specific, strong intent", 90),
        ("#aitutorial", 10, "Educational, high save rate", 75),
        ("#artificialintelligence", 9, "Professional audience skew", 70),
        ("#aitools", 8, "Purchase intent, listicle magnet", 85),
    ]

    # Visual bar chart style
    for i, (tag, count, note, bar_pct) in enumerate(hashtags):
        y = Inches(1.4) + i * Inches(1.15)

        # Tag name
        add_text_box(slide, Inches(0.6), y, Inches(3), Inches(0.5),
                     tag, 24, WHITE, True)

        # Bar background
        bar_bg = add_shape(slide, Inches(3.8), y + Inches(0.05), Inches(7), Inches(0.4), DARK_GRAY, 0.15)

        # Bar fill
        bar_w = Inches(7 * bar_pct / 100)
        colors = [ACCENT_PINK, ACCENT_CYAN, GREEN, ACCENT_PURPLE, RGBColor(0xFF, 0xD7, 0x00)]
        bar_fill = add_shape(slide, Inches(3.8), y + Inches(0.05), bar_w, Inches(0.4), colors[i], 0.15)

        # Count on bar
        add_text_box(slide, Inches(3.8) + bar_w - Inches(1), y + Inches(0.08), Inches(0.9), Inches(0.35),
                     f"{count} videos", 12, WHITE, True, PP_ALIGN.RIGHT)

        # Note
        add_text_box(slide, Inches(3.8), y + Inches(0.5), Inches(7), Inches(0.4),
                     note, 12, GRAY, False)


def main():
    print("Loading data...")
    with open(DATA_FILE) as f:
        videos = json.load(f)

    # Sort by views
    videos.sort(key=lambda v: v.get("playCount", 0), reverse=True)

    print(f"Building presentation with {len(videos)} videos...")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    print("  Slide 1: Title")
    build_title_slide(prs, videos)

    # Slide 2: Categories overview
    print("  Slide 2: Categories")
    build_categories_slide(prs)

    # Slide 3: Top 5 videos with thumbnails
    print("  Slide 3: Top 5 videos")
    build_top_videos_slide(prs, videos, 0, 5, "TOP 5 VIRAL VIDEOS")

    # Slide 4: Videos 6-10
    print("  Slide 4: Videos 6-10")
    build_top_videos_slide(prs, videos, 5, 5, "VIRAL VIDEOS #6-10")

    # Slide 5: Videos 11-15
    print("  Slide 5: Videos 11-15")
    build_top_videos_slide(prs, videos, 10, 5, "VIRAL VIDEOS #11-15")

    # Slide 6: Top hooks
    print("  Slide 6: Top hooks")
    build_hooks_slide(prs)

    # Slide 7: Hashtag performance
    print("  Slide 7: Hashtags")
    build_hashtags_slide(prs)

    # Slide 8: Most mentioned tools
    print("  Slide 8: Tools")
    build_tools_slide(prs)

    # Slide 9: Takeaways
    print("  Slide 9: Takeaways")
    build_takeaways_slide(prs)

    print(f"\nSaving to {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    print("Done!")


if __name__ == "__main__":
    main()
